"""
PowerPoint Helper Utility

This module provides enhanced PowerPoint manipulation capabilities
for creating professional presentations with consistent styling.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from typing import Dict, Any, List, Optional, Tuple, Union
import io
import tempfile
import os
from pathlib import Path


class PPTXHelper:
    """
    Enhanced PowerPoint helper with professional styling and layout management.
    """
    
    def __init__(self):
        self.default_font = 'Calibri'
        self.title_font_size = Pt(24)
        self.subtitle_font_size = Pt(18)
        self.body_font_size = Pt(14)
        self.caption_font_size = Pt(10)
        
        # Color scheme (modern blue theme)
        self.colors = {
            'primary': RGBColor(0, 119, 181),      # Blue
            'secondary': RGBColor(0, 156, 218),     # Light Blue
            'accent': RGBColor(255, 192, 0),        # Gold
            'text': RGBColor(68, 68, 68),           # Dark Gray
            'background': RGBColor(248, 249, 250),   # Light Gray
            'white': RGBColor(255, 255, 255)
        }
    
    def create_presentation_from_template(self, template_path: Optional[str] = None) -> Presentation:
        """Create a new presentation from template or default."""
        if template_path and os.path.exists(template_path):
            try:
                return Presentation(template_path)
            except Exception as e:
                print(f"Error loading template {template_path}: {e}")
                return Presentation()
        else:
            return Presentation()
    
    def add_title_slide(
        self, 
        prs: Presentation, 
        title: str, 
        subtitle: str = "",
        layout_index: int = 0
    ) -> None:
        """Add a professional title slide."""
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        
        # Configure title
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = title
            self._style_text_shape(title_shape, font_size=self.title_font_size, 
                                 color=self.colors['primary'], bold=True, alignment=PP_ALIGN.CENTER)
        
        # Configure subtitle
        if len(slide.placeholders) > 1 and subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            self._style_text_shape(subtitle_shape, font_size=self.subtitle_font_size,
                                 color=self.colors['text'], alignment=PP_ALIGN.CENTER)
    
    def add_content_slide(
        self, 
        prs: Presentation, 
        title: str, 
        content: Union[str, List[str]],
        layout_index: int = 1
    ) -> None:
        """Add a content slide with bullet points or text."""
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        
        # Configure title
        title_shape = slide.shapes.title
        title_shape.text = title
        self._style_text_shape(title_shape, font_size=self.subtitle_font_size,
                             color=self.colors['primary'], bold=True)
        
        # Configure content
        if len(slide.placeholders) > 1:
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()
            
            if isinstance(content, str):
                # Single text block
                p = text_frame.paragraphs[0]
                p.text = content
                self._style_paragraph(p, font_size=self.body_font_size, color=self.colors['text'])
            else:
                # List of bullet points
                for i, item in enumerate(content):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = item
                    p.level = 0
                    self._style_paragraph(p, font_size=self.body_font_size, color=self.colors['text'])
    
    def add_image_slide(
        self, 
        prs: Presentation, 
        title: str, 
        image_buffer: io.BytesIO,
        description: str = "",
        layout_index: int = 5
    ) -> None:
        """Add a slide with an image and optional description."""
        slide_layout = prs.slide_layouts[layout_index]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_shape.text_frame
        title_frame.text = title
        self._style_text_shape(title_shape, font_size=self.subtitle_font_size,
                             color=self.colors['primary'], bold=True, alignment=PP_ALIGN.CENTER)
        
        # Add image
        self._add_image_from_buffer(slide, image_buffer, 
                                  left=Inches(0.5), top=Inches(1.3), 
                                  width=Inches(9), height=Inches(5.5))
        
        # Add description if provided
        if description:
            desc_shape = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.8))
            desc_frame = desc_shape.text_frame
            desc_frame.text = description
            self._style_text_shape(desc_shape, font_size=self.caption_font_size,
                                 color=self.colors['text'], alignment=PP_ALIGN.CENTER)
    
    def add_two_column_slide(
        self, 
        prs: Presentation, 
        title: str,
        left_content: Union[str, List[str], io.BytesIO],
        right_content: Union[str, List[str], io.BytesIO],
        layout_index: int = 3
    ) -> None:
        """Add a two-column layout slide."""
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        
        # Configure title
        title_shape = slide.shapes.title
        title_shape.text = title
        self._style_text_shape(title_shape, font_size=self.subtitle_font_size,
                             color=self.colors['primary'], bold=True)
        
        # Add left content
        left_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
        self._add_content_to_shape(left_shape, left_content)
        
        # Add right content
        right_shape = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4.5), Inches(5.5))
        self._add_content_to_shape(right_shape, right_content)
    
    def add_data_summary_slide(
        self, 
        prs: Presentation, 
        title: str,
        summary_data: Dict[str, Any]
    ) -> None:
        """Add a slide with data summary in a structured format."""
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Configure title
        title_shape = slide.shapes.title
        title_shape.text = title
        self._style_text_shape(title_shape, font_size=self.subtitle_font_size,
                             color=self.colors['primary'], bold=True)
        
        # Create content layout
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        # Add summary items
        for key, value in summary_data.items():
            p = text_frame.add_paragraph() if text_frame.paragraphs else text_frame.paragraphs[0]
            
            # Format key-value pairs nicely
            if isinstance(value, (int, float)):
                if isinstance(value, float):
                    formatted_value = f"{value:,.2f}" if value > 1 else f"{value:.4f}"
                else:
                    formatted_value = f"{value:,}"
            else:
                formatted_value = str(value)
            
            p.text = f"{key.replace('_', ' ').title()}: {formatted_value}"
            p.level = 0
            self._style_paragraph(p, font_size=self.body_font_size, color=self.colors['text'])
    
    def add_insights_slide(
        self, 
        prs: Presentation, 
        title: str,
        insights: List[str],
        recommendations: List[str] = None
    ) -> None:
        """Add a slide with insights and recommendations."""
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Configure title
        title_shape = slide.shapes.title
        title_shape.text = title
        self._style_text_shape(title_shape, font_size=self.subtitle_font_size,
                             color=self.colors['primary'], bold=True)
        
        # Create content
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        # Add insights section
        if insights:
            # Section header
            p = text_frame.paragraphs[0]
            p.text = "Key Insights:"
            self._style_paragraph(p, font_size=Pt(16), color=self.colors['secondary'], bold=True)
            
            # Insights list
            for insight in insights[:5]:  # Limit to 5 insights
                p = text_frame.add_paragraph()
                p.text = f"• {insight}"
                p.level = 0
                self._style_paragraph(p, font_size=self.body_font_size, color=self.colors['text'])
        
        # Add recommendations section
        if recommendations:
            # Add spacing
            text_frame.add_paragraph()
            
            # Section header
            p = text_frame.add_paragraph()
            p.text = "Recommendations:"
            self._style_paragraph(p, font_size=Pt(16), color=self.colors['secondary'], bold=True)
            
            # Recommendations list
            for rec in recommendations[:3]:  # Limit to 3 recommendations
                p = text_frame.add_paragraph()
                p.text = f"• {rec}"
                p.level = 0
                self._style_paragraph(p, font_size=self.body_font_size, color=self.colors['text'])
    
    def add_metrics_dashboard_slide(
        self, 
        prs: Presentation, 
        title: str,
        metrics: Dict[str, Union[int, float, str]]
    ) -> None:
        """Add a slide with key metrics displayed as a dashboard."""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_shape.text_frame
        title_frame.text = title
        self._style_text_shape(title_shape, font_size=self.subtitle_font_size,
                             color=self.colors['primary'], bold=True, alignment=PP_ALIGN.CENTER)
        
        # Create metric boxes
        metrics_list = list(metrics.items())
        cols = 3
        rows = (len(metrics_list) + cols - 1) // cols
        
        box_width = Inches(2.8)
        box_height = Inches(1.2)
        start_left = Inches(0.8)
        start_top = Inches(1.8)
        
        for i, (metric_name, metric_value) in enumerate(metrics_list[:9]):  # Max 9 metrics
            row = i // cols
            col = i % cols
            
            left = start_left + col * (box_width + Inches(0.3))
            top = start_top + row * (box_height + Inches(0.4))
            
            # Create metric box
            self._create_metric_box(slide, left, top, box_width, box_height, 
                                  metric_name, metric_value)
    
    def _create_metric_box(
        self, 
        slide, 
        left: Inches, 
        top: Inches, 
        width: Inches, 
        height: Inches,
        metric_name: str, 
        metric_value: Union[int, float, str]
    ) -> None:
        """Create a styled metric box on the slide."""
        # Background shape
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = self.colors['background']
        bg_shape.line.color.rgb = self.colors['secondary']
        bg_shape.line.width = Pt(2)
        
        # Metric value (large text)
        value_shape = slide.shapes.add_textbox(
            left + Inches(0.1), top + Inches(0.1), 
            width - Inches(0.2), height * 0.6
        )
        value_frame = value_shape.text_frame
        value_frame.margin_left = value_frame.margin_right = 0
        value_frame.margin_top = value_frame.margin_bottom = 0
        
        # Format value
        if isinstance(metric_value, float):
            formatted_value = f"{metric_value:,.2f}" if metric_value > 1 else f"{metric_value:.4f}"
        elif isinstance(metric_value, int):
            formatted_value = f"{metric_value:,}"
        else:
            formatted_value = str(metric_value)
        
        value_frame.text = formatted_value
        self._style_text_shape(value_shape, font_size=Pt(20), color=self.colors['primary'],
                             bold=True, alignment=PP_ALIGN.CENTER)
        
        # Metric name (smaller text)
        name_shape = slide.shapes.add_textbox(
            left + Inches(0.1), top + height * 0.6, 
            width - Inches(0.2), height * 0.4
        )
        name_frame = name_shape.text_frame
        name_frame.margin_left = name_frame.margin_right = 0
        name_frame.margin_top = name_frame.margin_bottom = 0
        name_frame.text = metric_name.replace('_', ' ').title()
        self._style_text_shape(name_shape, font_size=Pt(10), color=self.colors['text'],
                             alignment=PP_ALIGN.CENTER)
    
    def _add_image_from_buffer(
        self, 
        slide, 
        image_buffer: io.BytesIO, 
        left: Inches, 
        top: Inches, 
        width: Inches, 
        height: Inches
    ) -> None:
        """Add image from buffer to slide."""
        # Save buffer to temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
            image_buffer.seek(0)
            tmp_file.write(image_buffer.getvalue())
            tmp_path = tmp_file.name
        
        try:
            slide.shapes.add_picture(tmp_path, left, top, width, height)
        finally:
            os.unlink(tmp_path)  # Clean up temp file
    
    def _add_content_to_shape(self, shape, content: Union[str, List[str], io.BytesIO]) -> None:
        """Add content to a text shape."""
        if isinstance(content, io.BytesIO):
            # Handle image content
            # Note: This would need additional handling for images in text boxes
            return
        
        text_frame = shape.text_frame
        text_frame.clear()
        
        if isinstance(content, str):
            # Single text block
            p = text_frame.paragraphs[0]
            p.text = content
            self._style_paragraph(p, font_size=self.body_font_size, color=self.colors['text'])
        else:
            # List of items
            for i, item in enumerate(content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = f"• {item}"
                p.level = 0
                self._style_paragraph(p, font_size=self.body_font_size, color=self.colors['text'])
    
    def _style_text_shape(
        self, 
        shape, 
        font_size: Pt = None, 
        color: RGBColor = None,
        bold: bool = False, 
        alignment: PP_ALIGN = None
    ) -> None:
        """Apply styling to a text shape."""
        text_frame = shape.text_frame
        
        for paragraph in text_frame.paragraphs:
            self._style_paragraph(paragraph, font_size, color, bold, alignment)
    
    def _style_paragraph(
        self, 
        paragraph, 
        font_size: Pt = None, 
        color: RGBColor = None,
        bold: bool = False, 
        alignment: PP_ALIGN = None
    ) -> None:
        """Apply styling to a paragraph."""
        if alignment:
            paragraph.alignment = alignment
        
        for run in paragraph.runs:
            font = run.font
            font.name = self.default_font
            
            if font_size:
                font.size = font_size
            if color:
                font.color.rgb = color
            if bold:
                font.bold = bold
    
    def apply_consistent_formatting(self, prs: Presentation) -> None:
        """Apply consistent formatting across all slides."""
        for slide in prs.slides:
            # Style title if present
            if hasattr(slide, 'shapes') and slide.shapes.title:
                self._style_text_shape(
                    slide.shapes.title,
                    font_size=self.subtitle_font_size,
                    color=self.colors['primary'],
                    bold=True
                )
            
            # Style other text shapes
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    # Skip if it's the title
                    if shape == slide.shapes.title:
                        continue
                    
                    self._style_text_shape(
                        shape,
                        font_size=self.body_font_size,
                        color=self.colors['text']
                    )
    
    def add_footer(self, prs: Presentation, footer_text: str) -> None:
        """Add footer to all slides."""
        for slide in prs.slides:
            footer_shape = slide.shapes.add_textbox(
                Inches(0.5), Inches(7.2), Inches(9), Inches(0.3)
            )
            footer_frame = footer_shape.text_frame
            footer_frame.text = footer_text
            self._style_text_shape(
                footer_shape,
                font_size=Pt(8),
                color=self.colors['text'],
                alignment=PP_ALIGN.CENTER
            )
    
    def save_presentation(self, prs: Presentation, output_path: str) -> str:
        """Save presentation with error handling."""
        try:
            # Ensure directory exists
            output_dir = Path(output_path).parent
            output_dir.mkdir(parents=True, exist_ok=True)
            
            # Apply final formatting
            self.apply_consistent_formatting(prs)
            
            # Save presentation
            prs.save(output_path)
            return f"Presentation saved successfully to: {output_path}"
        
        except Exception as e:
            return f"Error saving presentation: {str(e)}"
