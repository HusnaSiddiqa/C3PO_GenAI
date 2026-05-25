"""
PowerPoint Generator Tool using LangChain

This module provides a comprehensive tool for generating PowerPoint presentations
from CSV data with automated visualizations and summaries using LLM capabilities.
"""

import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from langchain.tools import BaseTool
from typing import Dict, Any, List, Optional, Union
import tempfile
import os
import io
import base64
from pathlib import Path

class PPTGenerator(BaseTool):
    """
    LangChain tool for generating PowerPoint presentations from CSV data and templates.
    """
    
    name: str = "ppt_generator"
    description: str = """
    A tool that takes a template PowerPoint file and CSV data to generate 
    a comprehensive presentation with visualizations and summaries.
    Accepts template_path, csv_path, output_path, and llm as parameters.
    """
    
    def __init__(self, llm=None, **kwargs):
        super().__init__(**kwargs)
        self._llm = llm
        
    def _run(
        self, 
        template_path: str, 
        df: pd.DataFrame, 
        output_path: str,
        llm: Optional[Any] = None,
        **kwargs
    ) -> str:
        """
        Generate PowerPoint presentation from template and CSV data.
        
        Args:
            template_path: Path to the template PPTX file
            csv_path: Path to the CSV data file
            output_path: Path where the generated PPTX will be saved
            llm: Language model for generating summaries and insights
            
        Returns:
            Success message with output path
        """
        try:

            model = llm or self._llm
            if not model:
                raise ValueError("LLM model is required for generating summaries")
            data_summary = kwargs.get('data_summary', {})
            visualizations = kwargs.get('visualizations', [])
            presentation = self._create_presentation(
                template_path, df, data_summary, visualizations, model
            )
            print("Presentation creating locally>>>>")
            presentation.save(output_path)
            
            return f"PowerPoint presentation successfully generated at: {output_path}"
            
        except Exception as e:
            print(f"Error generating PowerPoint: {e}")
            return f"Error generating PowerPoint: {str(e)}"
    
    def _load_csv_data(self, csv_path: str) -> pd.DataFrame:
        """Load and validate CSV data."""
        if not os.path.exists(csv_path):
            raise FileNotFoundError(f"CSV file not found: {csv_path}")
        
        df = pd.read_csv(csv_path)
        if df.empty:
            raise ValueError("CSV file is empty")
        
        return df
    
    def _analyze_data(self, df: pd.DataFrame, llm: Any) -> Dict[str, Any]:
        """Analyze data and generate summary using LLM."""
        # Basic data statistics
        stats = {
            'rows': len(df),
            'columns': len(df.columns),
            'column_names': list(df.columns),
            'data_types': df.dtypes.to_dict(),
            'missing_values': df.isnull().sum().to_dict(),
            'numeric_summary': df.describe().to_dict() if len(df.select_dtypes(include=['number']).columns) > 0 else {}
        }
        
        # Generate LLM-based insights
        data_preview = df.head(10).to_string()
        data_info = f"""
        Dataset Information:
        - Rows: {stats['rows']}
        - Columns: {stats['columns']}
        - Column Names: {', '.join(stats['column_names'])}
        
        Sample Data:
        {data_preview}
        
        Data Types:
        {stats['data_types']}
        """
        
        prompt = f"""
        Analyze the following dataset and provide:
        1. A brief executive summary (2-3 sentences)
        2. Key insights and patterns
        3. Recommendations for visualizations
        4. Important trends or anomalies
        
        Dataset Information:
        {data_info}
        
        Please provide a concise but comprehensive analysis.
        """
        
        try:
            llm_analysis = llm.invoke(prompt)
            if hasattr(llm_analysis, 'content'):
                analysis_text = llm_analysis.content
            else:
                analysis_text = str(llm_analysis)
        except Exception as e:
            analysis_text = f"LLM analysis unavailable: {str(e)}"
        
        stats['llm_analysis'] = analysis_text
        return stats
    
    def _create_visualizations(self, df: pd.DataFrame, llm: Any) -> List[Dict[str, Any]]:
        """Create appropriate visualizations based on data."""
        visualizations = []
        
        # Set style for better looking plots
        plt.style.use('seaborn-v0_8')
        sns.set_palette("husl")
        
        numeric_cols = df.select_dtypes(include=['number']).columns
        categorical_cols = df.select_dtypes(include=['object']).columns
        
        # Create correlation heatmap if multiple numeric columns
        if len(numeric_cols) > 1:
            vis_data = self._create_correlation_heatmap(df[numeric_cols])
            if vis_data:
                visualizations.append(vis_data)
        
        # Create distribution plots for numeric columns
        for col in numeric_cols[:3]:  # Limit to first 3 numeric columns
            vis_data = self._create_distribution_plot(df, col)
            if vis_data:
                visualizations.append(vis_data)
        
        # Create bar charts for categorical columns
        for col in categorical_cols[:2]:  # Limit to first 2 categorical columns
            if df[col].nunique() <= 20:  # Only if reasonable number of categories
                vis_data = self._create_bar_chart(df, col)
                if vis_data:
                    visualizations.append(vis_data)
        
        # Create scatter plot if we have at least 2 numeric columns
        if len(numeric_cols) >= 2:
            vis_data = self._create_scatter_plot(df, numeric_cols[0], numeric_cols[1])
            if vis_data:
                visualizations.append(vis_data)
        
        return visualizations
    
    def _create_correlation_heatmap(self, numeric_df: pd.DataFrame) -> Optional[Dict[str, Any]]:
        """Create correlation heatmap."""
        try:
            fig, ax = plt.subplots(figsize=(10, 8))
            correlation_matrix = numeric_df.corr()
            sns.heatmap(correlation_matrix, annot=True, cmap='coolwarm', center=0, ax=ax)
            ax.set_title('Correlation Matrix', fontsize=16, fontweight='bold')
            
            # Save to bytes
            img_buffer = io.BytesIO()
            plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight')
            img_buffer.seek(0)
            plt.close()
            
            return {
                'type': 'heatmap',
                'title': 'Correlation Analysis',
                'image_buffer': img_buffer,
                'description': 'Correlation matrix showing relationships between numeric variables'
            }
        except Exception as e:
            print(f"Error creating correlation heatmap: {e}")
            return None
    
    def _create_distribution_plot(self, df: pd.DataFrame, column: str) -> Optional[Dict[str, Any]]:
        """Create distribution plot for a numeric column."""
        try:
            fig, ax = plt.subplots(figsize=(10, 6))
            df[column].hist(bins=30, alpha=0.7, ax=ax)
            ax.set_title(f'Distribution of {column}', fontsize=14, fontweight='bold')
            ax.set_xlabel(column)
            ax.set_ylabel('Frequency')
            ax.grid(True, alpha=0.3)
            
            # Save to bytes
            img_buffer = io.BytesIO()
            plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight')
            img_buffer.seek(0)
            plt.close()
            
            return {
                'type': 'histogram',
                'title': f'Distribution of {column}',
                'image_buffer': img_buffer,
                'description': f'Frequency distribution of {column} values'
            }
        except Exception as e:
            print(f"Error creating distribution plot for {column}: {e}")
            return None
    
    def _create_bar_chart(self, df: pd.DataFrame, column: str) -> Optional[Dict[str, Any]]:
        """Create bar chart for categorical column."""
        try:
            fig, ax = plt.subplots(figsize=(12, 6))
            value_counts = df[column].value_counts().head(10)
            value_counts.plot(kind='bar', ax=ax)
            ax.set_title(f'Top Values in {column}', fontsize=14, fontweight='bold')
            ax.set_xlabel(column)
            ax.set_ylabel('Count')
            ax.tick_params(axis='x', rotation=45)
            
            # Save to bytes
            img_buffer = io.BytesIO()
            plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight')
            img_buffer.seek(0)
            plt.close()
            
            return {
                'type': 'bar_chart',
                'title': f'Top Values in {column}',
                'image_buffer': img_buffer,
                'description': f'Bar chart showing frequency of top values in {column}'
            }
        except Exception as e:
            print(f"Error creating bar chart for {column}: {e}")
            return None
    
    def _create_scatter_plot(self, df: pd.DataFrame, x_col: str, y_col: str) -> Optional[Dict[str, Any]]:
        """Create scatter plot between two numeric columns."""
        try:
            fig, ax = plt.subplots(figsize=(10, 8))
            df.plot.scatter(x=x_col, y=y_col, alpha=0.6, ax=ax)
            ax.set_title(f'{x_col} vs {y_col}', fontsize=14, fontweight='bold')
            ax.grid(True, alpha=0.3)
            
            # Save to bytes
            img_buffer = io.BytesIO()
            plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight')
            img_buffer.seek(0)
            plt.close()
            
            return {
                'type': 'scatter',
                'title': f'{x_col} vs {y_col}',
                'image_buffer': img_buffer,
                'description': f'Scatter plot showing relationship between {x_col} and {y_col}'
            }
        except Exception as e:
            print(f"Error creating scatter plot: {e}")
            return None
    
    def _create_presentation(
        self, 
        template_path: str, 
        df: pd.DataFrame, 
        data_summary: Dict[str, Any], 
        visualizations: List[Dict[str, Any]],
        llm: Any
    ) -> Presentation:
        """Create the PowerPoint presentation."""
        # Load template or create new presentation
        if os.path.exists(template_path):
            prs = Presentation(template_path)
        else:
            prs = Presentation()
            print(f"Template not found at {template_path}, creating new presentation")
        # print("Added title slide")
        # # Add title slide
        # self._add_title_slide(prs, data_summary)
        # print("Added data overview slide")
        # # Add data overview slide
        # self._add_data_overview_slide(prs, data_summary)
        print("Added visualization slides")
        # Add visualization slides
        for viz in visualizations:
            self._add_visualization_slide(prs, viz)
        print("Added insights and summary slide")
        # Add insights and summary slide
        self._add_insights_slide(prs, data_summary, llm)
        
        return prs
    
    def _add_title_slide(self, prs: Presentation, data_summary: Dict[str, Any]):
        """Add title slide to presentation."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Data Analysis Report"
        subtitle.text = f"Analysis of {data_summary['rows']} records across {data_summary['columns']} variables"
    
    def _add_data_overview_slide(self, prs: Presentation, data_summary: Dict[str, Any]):
        """Add data overview slide."""
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Data Overview"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = f"Dataset Summary:\n"
        
        p = tf.add_paragraph()
        p.text = f"• Total Records: {data_summary['rows']:,}"
        
        p = tf.add_paragraph()
        p.text = f"• Number of Variables: {data_summary['columns']}"
        
        p = tf.add_paragraph()
        p.text = f"• Variables: {', '.join(data_summary['column_names'][:5])}"
        if len(data_summary['column_names']) > 5:
            p.text += "..."
        
        # Add missing values info if any
        missing_info = data_summary.get('missing_values', {})
        if any(missing_info.values()):
            p = tf.add_paragraph()
            p.text = "• Data Quality: Some missing values detected"
    
    def _add_visualization_slide(self, prs: Presentation, viz_data: Dict[str, Any]):
        """Add visualization slide."""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = viz_data['title']
        title_frame.paragraphs[0].font.size = Inches(0.3)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add image
        img_buffer = viz_data['image_buffer']
        img_buffer.seek(0)
        
        # Save buffer to temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
            tmp_file.write(img_buffer.getvalue())
            tmp_path = tmp_file.name
        
        try:
            slide.shapes.add_picture(tmp_path, Inches(1), Inches(1.5), Inches(8), Inches(5))
        finally:
            os.unlink(tmp_path)  # Clean up temp file
        
        # Add description
        desc_shape = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.7))
        desc_frame = desc_shape.text_frame
        desc_frame.text = viz_data['description']
        desc_frame.paragraphs[0].font.size = Inches(0.15)
    
    def _add_insights_slide(self, prs: Presentation, data_summary: Dict[str, Any], llm: Any):
        """Add insights and summary slide."""
        if data_summary and 'unavailable' not in data_summary.lower():
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            title.text = "Key Insights & Recommendations"
            
            content = slide.placeholders[1]
            tf = content.text_frame
            
            # Add LLM analysis if available
            llm_analysis = data_summary
            tf.text = llm_analysis
        

    async def _arun(self, *args, **kwargs):
        """Async version of the tool (not implemented)."""
        raise NotImplementedError("Async version not implemented")


def create_ppt_from_csv(
    template_path: str,
    csv_path: str,
    output_path: str,
    llm: Any
) -> str:
    """
    Convenience function to generate PowerPoint from CSV data.
    
    Args:
        template_path: Path to the template PPTX file
        csv_path: Path to the CSV data file
        output_path: Path where the generated PPTX will be saved
        llm: Language model for generating summaries and insights
        
    Returns:
        Success message with output path
    """
    generator = PPTGenerator(llm=llm)
    return generator._run(template_path, csv_path, output_path, llm)
