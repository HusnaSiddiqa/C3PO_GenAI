import os
import boto3
import requests
import numpy as np
import traceback
from io import BytesIO
from typing import List, Dict, Optional
from dataclasses import dataclass, field

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_aws import BedrockEmbeddings
from langchain_core.tools import tool
from botocore.exceptions import ClientError
from rank_bm25 import BM25Okapi


import fitz
import docx
import json
import csv
from pptx import Presentation
from openpyxl import load_workbook


@dataclass
class DocumentChunk:
    """Represents a single chunk of a document with its content, embedding, and metadata."""
    content: str
    embedding: np.ndarray
    chunk_id: str
    metadata: Dict[str, any]


@dataclass
class ProcessedDocument:
    """Holds metadata and chunk/embedding information for a processed document."""
    document_path: str
    total_chunks: int
    chunks: List[DocumentChunk]
    bm25_index: 'BM25Okapi' = field(repr=False)
    chunk_embeddings: np.ndarray = field(repr=False)
    processing_metadata: Dict[str, any]



class DocumentProcessor:
    """Handles loading, parsing, and chunking documents from S3 by character count."""

    def __init__(self, embedding_model: BedrockEmbeddings):
        self.embedding_model = embedding_model
        self.s3_client = boto3.client("s3")

    def load_document_from_s3(self, s3_path: str) -> str:
        """Fetches a file from S3 and extracts its text."""
        if not s3_path.startswith("s3://"):
            raise ValueError("Invalid S3 path format.")

        bucket, key = s3_path[5:].split("/", 1)
        try:
            response = self.s3_client.get_object(Bucket=bucket, Key=key)
            extracted_text = self._extract_text(response["Body"].read(), s3_path)
            return extracted_text
        except ClientError as e:
            error_code = e.response['Error']['Code']
            error_message = e.response['Error']['Message']
            return f"❌ AWS Error: {error_code}"
        except Exception as e:
            print(traceback.format_exc())
            return f"❌ General Error: {e}"

    def _extract_text(self, content_bytes: bytes, file_path: str) -> str:
        """Extracts text from supported file types."""
        ext = os.path.splitext(file_path)[1].lower()
        try:
            if ext == ".pdf":
                with fitz.open(stream=content_bytes, filetype="pdf") as doc:
                    return "\n".join(page.get_text() for page in doc)
            elif ext == ".docx":
                doc = docx.Document(BytesIO(content_bytes))
                return "\n".join(p.text for p in doc.paragraphs)
            elif ext in [".ppt", ".pptx"]:
                prs = Presentation(BytesIO(content_bytes))
                return "\n".join(
                    run.text for slide in prs.slides for shape in slide.shapes
                    if shape.has_text_frame for p in shape.text_frame.paragraphs
                    for run in p.runs if run.text.strip()
                )
            elif ext == ".xlsx":
                wb = load_workbook(BytesIO(content_bytes), read_only=True)
                return "\n".join(
                    ", ".join(str(cell.value) for cell in row if cell.value)
                    for sheet in wb.worksheets for row in sheet.iter_rows()
                )
            elif ext == ".csv":
                lines = content_bytes.decode("utf-8", errors="ignore").splitlines()
                return "\n".join([", ".join(row) for row in csv.reader(lines)])
            elif ext == ".json":
                data = json.loads(content_bytes.decode("utf-8", errors="ignore"))
                return json.dumps(data, indent=2)
            elif ext == ".txt":
                return content_bytes.decode("utf-8", errors="ignore")
            else:
                return f"File type {ext} not supported."
        except Exception as e:
            return f"Error parsing document: {e}"

    async def create_embeddings(self, chunks: List[str]) -> List[np.ndarray]:
        """Converts a list of text chunks into embeddings, handling batching for large lists."""
        all_embeddings = []
        batch_size = 500
        for i in range(0, len(chunks), batch_size):
            batch_chunks = chunks[i:i + batch_size]
            try:
                batch_embeddings = await self.embedding_model.aembed_documents(batch_chunks)
                all_embeddings.extend(batch_embeddings)
            except Exception as e:
                all_embeddings.extend([[0.0] * 1536 for _ in batch_chunks])  # Assuming embedding size
        return [np.array(emb) for emb in all_embeddings]

    async def process_document(self, document_path: str) -> Optional[ProcessedDocument]:
        """Orchestrates the document processing: load, split, embed, index chunks."""
        try:
            full_text = self.load_document_from_s3(document_path)
            if full_text.startswith("❌"):
                return None

            text_splitter = RecursiveCharacterTextSplitter(chunk_size=300, chunk_overlap=30, length_function=len)
            chunks = text_splitter.split_text(full_text)
            if not chunks: return None

            embeddings = await self.create_embeddings(chunks)
            if not embeddings: return None

            bm25_index = BM25Okapi([c.split() for c in chunks])
            doc_chunks = [
                DocumentChunk(
                    content=text, embedding=emb, chunk_id=f"{os.path.basename(document_path)}_chunk_{i:04d}",
                    metadata={"chunk_index": i}
                ) for i, (text, emb) in enumerate(zip(chunks, embeddings))
            ]
            return ProcessedDocument(
                document_path=document_path, total_chunks=len(doc_chunks), chunks=doc_chunks,
                bm25_index=bm25_index, chunk_embeddings=np.array(embeddings),
                processing_metadata={"original_char_count": len(full_text)}
            )
        except Exception as e:
            print(f"❌ Document processing failed: {e}")
            return None


# --- Retriever & Tool Definition (Slightly Refined) ---

class EnhancedRetriever:
    """Handles caching, processing, and hybrid retrieval of documents."""

    def __init__(self, embedding_model: BedrockEmbeddings):
        self.embedding_model = embedding_model
        self.processor = DocumentProcessor(embedding_model)
        self.processed_documents: Dict[str, ProcessedDocument] = {}

    async def _ensure_document_is_processed(self, document_path: str) -> bool:
        """Processes a document if it's not already in the cache."""
        if document_path in self.processed_documents:
            print(f"[DEBUG] Document '{document_path}' found in cache.")
            return True
        processed_doc = await self.processor.process_document(document_path)
        if processed_doc:
            self.processed_documents[document_path] = processed_doc
            return True
        else:
            print(f"❌ Failed to process and cache document: {document_path}")
            return False

    def _cosine_similarity(self, a: np.ndarray, b: np.ndarray) -> float:
        return np.dot(a, b) / (np.linalg.norm(a) * np.linalg.norm(b) + 1e-10)

    async def retrieve_relevant_chunks(self, query: str, document_path: str, top_k: int = 50) -> List[str]:
        """Performs the hybrid search for a given query and document."""

        if not await self._ensure_document_is_processed(document_path):
            return ["❌ Failed to process document."]

        doc = self.processed_documents[document_path]
        query_emb = await self.embedding_model.aembed_query(query)
        bm25_scores = doc.bm25_index.get_scores(query.split())
        semantic_scores = [self._cosine_similarity(query_emb, emb) for emb in doc.chunk_embeddings]

        rrf_scores = {}
        k = 60
        for rank, (i, _) in enumerate(sorted(enumerate(bm25_scores), key=lambda x: x[1], reverse=True)):
            rrf_scores[i] = rrf_scores.get(i, 0) + 1 / (k + rank + 1)
        for rank, (i, _) in enumerate(sorted(enumerate(semantic_scores), key=lambda x: x[1], reverse=True)):
            rrf_scores[i] = rrf_scores.get(i, 0) + 1 / (k + rank + 1)

        final_indices = [i for i, _ in sorted(rrf_scores.items(), key=lambda x: x[1], reverse=True)[:top_k]]
        retrieved_chunks = [doc.chunks[i].content for i in final_indices]
        return retrieved_chunks



_retriever_instance: Optional[EnhancedRetriever] = None


def set_retriever_instance(r: EnhancedRetriever):
    global _retriever_instance
    _retriever_instance = r


@tool
async def retrieve_document_content(query: str, document_path: str) -> str:
    """Retrieves relevant content from a document based on the user's query."""
    if _retriever_instance is None:
        return "❌ Retriever not initialized."

    chunks = await _retriever_instance.retrieve_relevant_chunks(query, document_path)
    if not chunks or "❌" in chunks[0]:
        return chunks[0] if chunks else "No relevant content found in the document."

    final_context = "\n\n---\n\n".join(chunks)
    print(final_context)
    return final_context

## Retrieve file content using conversation_id using chat_manager api
def retrieve_file_content(conversation_id: str) -> str:
    """API call Retrieves the content of a file based on the conversation ID."""
    response = requests.get(f"http://{os.getenv('CHAT_MANAGER_BASE_URL')}:{os.getenv('CHAT_MANAGER_PORT')}/v2/chat-manager/chat/file/{conversation_id}")
    print("===========response=======", response)
    return response.json()

class Tools:
    retrieve_document_content = retrieve_document_content